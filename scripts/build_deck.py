"""Generate the presentation deck as a real, editable .pptx.

Mirrors docs/index.html in running order and palette, so the slides and the
landing page tell one story. Numbers here are the real ones (91 tests, the live
model and spend captured in the last snapshot); update them if they move.

    .venv/bin/python scripts/build_deck.py   ->  Renewal-Risk-Platform.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent

# Palette lifted from docs/index.html so the deck and the page match.
INK    = RGBColor(0x1A, 0x1A, 0x18)
PAPER  = RGBColor(0xFB, 0xFB, 0xFA)
PANEL  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x6B, 0x6B, 0x66)
LINE   = RGBColor(0xE5, 0xE4, 0xE0)
ACCENT = RGBColor(0xB8, 0x49, 0x2F)
OK     = RGBColor(0x2F, 0x7D, 0x55)
CODEBG = RGBColor(0xF4, 0xF3, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Arial"
MONO   = "Menlo"

EMU = Inches(13.333)


def _p(tf):
    first = tf.paragraphs[0]
    if not first.runs and (first.text or "") == "":
        return first
    return tf.add_paragraph()


def para(tf, text, size=18, color=INK, bold=False, italic=False, font=FONT,
         after=8, before=0, align=PP_ALIGN.LEFT, level=0):
    p = _p(tf)
    p.alignment = align
    p.level = level
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    run = p.add_run()
    run.text = text
    f = run.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    return p


def runs(tf, parts, size=18, after=8, align=PP_ALIGN.LEFT):
    """One paragraph, several runs with different styling: [(text, {opts})]."""
    p = _p(tf)
    p.alignment = align
    p.space_after = Pt(after)
    for text, opt in parts:
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", False)
        f.italic = opt.get("italic", False)
        f.name = opt.get("font", FONT)
        f.color.rgb = opt.get("color", INK)
    return p


def box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def rect(slide, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def blank(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def header(slide, eyebrow, heading, heading_size=32):
    tf = box(slide, 0.7, 0.52, 12, 0.4)
    para(tf, eyebrow.upper(), size=12.5, color=ACCENT, bold=True, after=0)
    tf2 = box(slide, 0.66, 0.92, 12, 1.2)
    para(tf2, heading, size=heading_size, color=INK, bold=True, after=0)
    rect(slide, 0.72, 1.74, 1.05, 0.045, ACCENT)


def footer(slide, n):
    tf = box(slide, 0.7, 7.02, 11, 0.35)
    runs(tf, [("Renewal Risk Platform", {"size": 9.5, "color": MUTED}),
              (f"   ·   {n}", {"size": 9.5, "color": MUTED})], after=0)


def say(slide, l, t, w, h, label, text):
    panel = rect(slide, l, t, w, h, PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    panel.adjustments[0] = 0.04
    rect(slide, l, t, 0.06, h, OK)
    tf = box(slide, l + 0.28, t + 0.16, w - 0.5, h - 0.3, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label.upper(), size=11, color=OK, bold=True, after=5)
    para(tf, text, size=14.5, color=INK, after=0)


def bullets(tf, items, size=17, after=10, color=INK):
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            runs(tf, [("→  ", {"color": ACCENT, "bold": True, "size": size}),
                      (lead, {"bold": True, "size": size, "color": color}),
                      (rest, {"size": size, "color": color})], after=after)
        else:
            runs(tf, [("→  ", {"color": ACCENT, "bold": True, "size": size}),
                      (it, {"size": size, "color": color})], after=after)


def stat(slide, l, t, w, big, label, big_color=INK):
    rect(slide, l, t, w, 1.5, PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = box(slide, l + 0.2, t + 0.22, w - 0.4, 1.1, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, size=27, color=big_color, bold=True, after=2)
    para(tf, label.upper(), size=10.5, color=MUTED, bold=True, after=0)


def table(slide, l, t, w, rows, col_w, header_row, body_rows, size=13):
    n = len(body_rows) + 1
    gt = slide.shapes.add_table(n, len(col_w), Inches(l), Inches(t),
                                Inches(w), Inches(0.5 * n)).table
    gt.first_row = False
    gt.horz_banding = False
    for i, cw in enumerate(col_w):
        gt.columns[i].width = Inches(cw)
    def cell(r, c, text, bold=False, color=INK, sz=size, mono=False):
        cl = gt.cell(r, c)
        cl.fill.solid()
        cl.fill.fore_color.rgb = PAPER if r == 0 else PANEL
        cl.margin_left = Inches(0.12)
        cl.margin_right = Inches(0.1)
        cl.margin_top = Inches(0.05)
        cl.margin_bottom = Inches(0.05)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cl.text_frame
        tf.word_wrap = True
        para(tf, text, size=sz, color=color, bold=bold,
             font=MONO if mono else FONT, after=0)
    for c, htext in enumerate(header_row):
        cell(0, c, htext.upper(), bold=True, color=MUTED, sz=10.5)
    for r, row in enumerate(body_rows, start=1):
        for c, (text, opt) in enumerate(row):
            cell(r, c, text, bold=opt.get("bold", False),
                 color=opt.get("color", INK), sz=opt.get("size", size),
                 mono=opt.get("mono", False))
    return gt


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = [0]

    def page():
        n[0] += 1
        return n[0]

    # 1 · Title (dark) ------------------------------------------------------
    s = blank(prs, INK)
    rect(s, 0.9, 2.5, 1.05, 0.05, ACCENT)
    tf = box(s, 0.85, 2.7, 11.6, 2.4)
    para(tf, "Renewal Risk Analyser & Router", size=44, color=WHITE, bold=True, after=10)
    para(tf, "One agent was the brief. The platform many agents plug into is the job.",
         size=20, color=RGBColor(0xC9, 0xC7, 0xC0), after=0)
    tf2 = box(s, 0.85, 5.6, 11.6, 1)
    runs(tf2, [("Denis Miano", {"size": 15, "color": WHITE, "bold": True}),
               ("     Supermetrics · AI Platform Engineer", {"size": 15, "color": RGBColor(0x98,0x95,0x8C)})],
         after=0)

    # 2 · Framing ----------------------------------------------------------
    s = blank(prs)
    header(s, "01 · The shape of it", "I built the platform, not a script")
    tf = box(s, 0.7, 2.1, 12, 1.2)
    para(tf, "The brief describes one agent. The role owns the layer many agents plug into. "
             "So the renewal agent is the first tenant on a shared platform, and two more run "
             "on the same rails without a line of the first changing.", size=18, color=MUTED, after=0)
    tf2 = box(s, 0.7, 3.5, 12, 3)
    para(tf2, "Three things it optimises for", size=15, color=ACCENT, bold=True, after=12)
    bullets(tf2, [
        ("Nothing the model says is trusted unchecked. ", "Every figure is verified against real data."),
        ("Every decision is explainable to a non-engineer. ", "Plain English, from the same trace."),
        ("When anything goes wrong, work routes to a human. ", "It never silently stops."),
    ], size=18, after=13)
    footer(s, page())

    # 3 · What it does -----------------------------------------------------
    s = blank(prs)
    header(s, "02 · What it does", "A customer is about to leave, and nobody noticed")
    tf = box(s, 0.7, 2.05, 6.15, 4.6)
    para(tf, "Health score drops as the renewal nears. The platform:", size=17, color=INK, after=12)
    bullets(tf, [
        "Gathers that account's facts from Salesforce, Gainsight, HubSpot, Zendesk",
        "Uses an LLM to work out the likely reason they might leave",
        "Writes the finding back to Salesforce and Gainsight",
        "Alerts the account owner in Slack with the numbers behind it",
        "Records every step, so anyone can ask why in plain English",
    ], size=16, after=11)
    # plain-english quote card
    rect(s, 7.15, 2.0, 5.5, 4.7, PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 7.15, 2.0, 0.06, 4.7, ACCENT)
    q = box(s, 7.45, 2.25, 5.05, 4.3)
    para(q, "Verdant Foods Co: high churn risk, because the person who championed us "
            "internally has moved on", size=15, color=INK, bold=True, after=10)
    for ln in [
        "1 · A renewal-approaching message arrived and health had fallen.",
        "2 · Looked the customer up and gathered 14 facts about them.",
        "3 · Judged the likely reason: the internal champion has moved on. 92% sure.",
        "4 · Rated it high priority: health 45, worth $67,000 a year.",
        "5 · Recorded it in Salesforce and Gainsight, alerted the owner in #cs-renewals.",
    ]:
        para(q, ln, size=12.5, color=MUTED, after=6)
    footer(s, page())

    # 4 · The hard part ----------------------------------------------------
    s = blank(prs)
    header(s, "03 · The hard part", "Three customers, near-identical triggers, three reasons")
    tf = box(s, 0.7, 2.0, 12, 0.8)
    para(tf, "All three lose health at a similar rate and renew soon. A system that just repeats "
             "the warning says the same useless thing about all three.", size=16, color=MUTED, after=0)
    table(s, 0.7, 2.95, 12, None,
          [3.0, 4.4, 4.6],
          ["Customer", "What the platform worked out", "What it means"],
          [
            [("Northwind  $84k", {"bold": True}),
             ("People largely stopped using what they pay for", {}),
             ("Team shrank, usage collapsed. Re-engage or lose them.", {"color": MUTED})],
            [("Bluefin  $31k", {"bold": True}),
             ("The data connections they depend on keep breaking", {}),
             ("Still daily users, but they distrust the numbers. Ours to fix.", {"color": MUTED})],
            [("Verdant  $67k", {"bold": True}),
             ("The internal champion has moved on", {}),
             ("Nothing broken, usage fine, but nobody left who knows it.", {"color": MUTED})],
          ], size=13.5)
    say(s, 0.7, 5.55, 12, 1.15, "The difference that matters",
        "Same trigger shape, same falling health, three different answers, each with its own "
        "evidence. That is an alert versus something worth acting on.")
    footer(s, page())

    # 5 · The finding ------------------------------------------------------
    s = blank(prs)
    header(s, "04 · The interesting part", "Your sample payload caught a real defect")
    tf = box(s, 0.7, 2.0, 12, 2.0)
    para(tf, "The prompt returned champion_loss for BOTH Verdant and Northwind, because "
             "Northwind's notes also mention a contact changing role. Two accounts, one answer: "
             "exactly the failure the file is built to expose.", size=17, color=INK, after=10)
    runs(tf, [("The discriminator was in the data, not the model. ", {"size": 17, "bold": True}),
              ("Northwind's automation also stopped (transfers paused, 18 of 22 seats idle): "
               "team-wide withdrawal. Verdant's automation is healthy and only the human work "
               "stopped: a lost champion.", {"size": 17, "color": INK})], after=0)
    table(s, 0.7, 4.35, 12, None, [3.4, 4.3, 4.3],
          ["Account", "First answer", "Now"],
          [
            [("Northwind Retail", {}), ("champion_loss", {"mono": True, "color": MUTED}),
             ("adoption_decline  ✓", {"mono": True, "color": OK})],
            [("Verdant Foods", {}), ("champion_loss", {"mono": True, "color": MUTED}),
             ("champion_loss  ✓", {"mono": True, "color": OK})],
          ], size=13.5)
    say(s, 0.7, 5.95, 12, 0.95, "Why I'd lead with this",
        "The prompt fix is not the point. A regression set caught it before a CS team did; "
        "all three are now permanent eval cases and the gate blocks any change that brings it back.")
    footer(s, page())

    # 6 · Architecture -----------------------------------------------------
    s = blank(prs)
    header(s, "05 · Architecture", "How I would run this in the cloud")
    img = ROOT / "docs" / "diagrams" / "cloud-architecture.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(8.55), Inches(1.15),
                             height=Inches(6.05))
    tf = box(s, 0.7, 2.15, 7.5, 4.6)
    bullets(tf, [
        ("Ingestion split from workers. ", "Vendors time out webhooks in seconds; the model "
         "takes tens. Shared, a slow model makes a vendor drop our events."),
        ("BigQuery, where the Golden Record lives. ", "Traces land where analysts join agent "
         "behaviour to revenue. “Did the accounts we flagged renew?” becomes SQL."),
        ("europe-north1. ", "EU data residency, closest region to Helsinki."),
    ], size=16, after=14)
    say(s, 0.7, 6.05, 7.5, 1.0, "At 10x, the first change",
        "Make ingestion asynchronous: validate, persist, enqueue, workers consume. That is the "
        "failure that bites first, before cost or storage.")
    footer(s, page())

    # 7 · Observability 1 --------------------------------------------------
    s = blank(prs)
    header(s, "06 · Observability", "“Why did this agent do that?” in fifteen seconds")
    tf = box(s, 0.7, 2.1, 5.9, 4.5)
    para(tf, "The product requirement, and the thing the panel cares about most.", size=17, color=MUTED, after=14)
    bullets(tf, [
        ("Plain English ", "for whoever is asking: no code, no identifiers, no JSON."),
        ("Rules and values ", "one click down for whoever is fixing."),
        ("One trace behind both, ", "so the two readings cannot disagree."),
        ("Skipped runs are traced too, ", "so “why did nothing happen?” also has an answer."),
    ], size=16, after=13)
    rect(s, 6.9, 2.05, 5.75, 4.75, INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c = box(s, 7.2, 2.35, 5.2, 4.2)
    para(c, "$ python cli.py why tr_fe83…", size=13, color=RGBColor(0xE0,0x7A,0x5F), font=MONO, after=10)
    for ln in [
        "Woken by 'renewal.approaching' for ACC-9033.",
        "Renewal in 24 days, health 52, down 14  (rule: act)",
        "Severity 'high': health 52 ≤ 55  (rule: band_2)",
        "Wrote to Salesforce + Gainsight  (rule: auto_approved)",
        "Routed to #cs-renewals  (rule: high_owner_alert)",
    ]:
        para(c, ln, size=12.5, color=RGBColor(0xEC,0xEA,0xE4), font=MONO, after=7)
    footer(s, page())

    # 8 · Observability 2 (new panels) -------------------------------------
    s = blank(prs)
    header(s, "07 · The platform watches itself", "Health, decisions, reviews and tests, on one page")
    para(box(s, 0.7, 1.95, 12, 0.6),
         "Legible to a CS lead, not just an engineer. Every item below is a live dashboard panel.",
         size=15.5, color=MUTED, after=0)
    cards = [
        ("Platform health check", "7 self-audit checks (owned, reviewed, nothing stuck, gate green, "
         "model answering, precision holding, clean). No LLM: each has a correct answer. Same checks gate CI."),
        ("Decision log", "Every rule every agent just fired, newest first, each click-through to its full reasoning."),
        ("Reviews & owners", "Per-agent last-checked and next-due date. Overdue agents are nudged to their owner in Slack."),
        ("Automated tests", "91 tests, collected live and bucketed by what they cover. A number, not a claim."),
    ]
    xs = [0.7, 6.75]
    ys = [2.65, 4.75]
    for i, (t, body) in enumerate(cards):
        l, top = xs[i % 2], ys[i // 2]
        rect(s, l, top, 5.9, 1.9, PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = box(s, l + 0.28, top + 0.2, 5.4, 1.55)
        para(tf, t, size=17, color=ACCENT, bold=True, after=6)
        para(tf, body, size=12.5, color=INK, after=0)
    footer(s, page())

    # 9 · Prove we used the key --------------------------------------------
    s = blank(prs)
    header(s, "08 · Proof, not assertion", "Did it really call the model, or fake it?")
    para(box(s, 0.7, 1.95, 12, 0.6),
         "The dashboard header shows LIVE or OFFLINE; the cost panel names the exact model, tokens and dollars.",
         size=16, color=MUTED, after=0)
    stat(s, 0.7, 2.85, 2.85, "LIVE", "run mode", big_color=OK)
    stat(s, 3.75, 2.85, 2.85, "20", "real model calls")
    stat(s, 6.8, 2.85, 2.85, "56.3k", "tokens in + out")
    stat(s, 9.85, 2.85, 2.75, "$0.31", "of the $50 grant")
    table(s, 0.7, 4.75, 12, None, [5.4, 2.0, 2.2, 2.4],
          ["Model actually used", "Calls", "Tokens", "Spend"],
          [
            [("anthropic/claude-sonnet-4.5", {"mono": True}), ("19", {}),
             ("53.9k", {}), ("$0.312", {})],
            [("openai/gpt-4o-mini  (chain fallback)", {"mono": True, "color": MUTED}), ("1", {}),
             ("2.4k", {}), ("$0.0004", {})],
          ], size=13.5)
    say(s, 0.7, 6.05, 12, 0.95, "The offline half of the proof",
        "Offline, nothing leaves the process and cost stays $0 with no model named. Which path "
        "ran is visible, not asserted, and the fallback hop above shows the model chain is real.")
    footer(s, page())

    # 10 · Bounded + fail to human -----------------------------------------
    s = blank(prs)
    header(s, "09 · When it can't be sure", "Bounded, and it fails to a human")
    tf = box(s, 0.7, 2.15, 5.9, 4.5)
    bullets(tf, [
        ("Retries are bounded. ", "Backoff on timeouts, never on a 4xx, capped attempts."),
        ("A circuit breaker ", "stops calling a vendor that is already down."),
        ("The daily budget caps spend ", "at 90%; runs degrade predictably, not mid-analysis."),
        ("Low confidence holds the write ", "and asks a person. It never guesses, never loops, never goes silent."),
    ], size=16.5, after=14)
    rect(s, 6.9, 2.15, 5.75, 4.3, PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c = box(s, 7.2, 2.45, 5.2, 3.8)
    para(c, "The four ways it degrades safely", size=13.5, color=ACCENT, bold=True, after=10)
    for ln in [
        "Invents evidence  →  discarded, rules used, flagged",
        "Model / vendor down  →  fallback, human still alerted",
        "Not confident  →  CRM write held, person asked",
        "Budget hit  →  stop spending, route to a human",
    ]:
        para(c, ln, size=13.5, color=INK, after=11)
    footer(s, page())

    # 11 · Adding an agent -------------------------------------------------
    s = blank(prs)
    header(s, "10 · The platform test", "Adding an agent is two files")
    tf = box(s, 0.7, 2.1, 6.0, 4.5)
    para(tf, "The real test of whether a platform exists: how little it takes to add the next one.",
         size=17, color=MUTED, after=14)
    bullets(tf, [
        ("One command ", "writes the module and the registry entry."),
        ("Refused ", "without an owner, a dotted event, or known tools."),
        ("Inherited free: ", "tracing, retries, idempotency, guardrails, permissions, the dashboard."),
        ("Least privilege: ", "it can touch only the tools its entry grants."),
    ], size=16, after=13)
    rect(s, 6.95, 2.1, 5.7, 2.5, INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c = box(s, 7.25, 2.4, 5.1, 2.0)
    para(c, "python cli.py new-agent expansion_signal \\", size=12.5, color=RGBColor(0xE0,0x7A,0x5F), font=MONO, after=5)
    para(c, "  --owner denis \\", size=12.5, color=RGBColor(0xEC,0xEA,0xE4), font=MONO, after=5)
    para(c, "  --subscribes-to health_score.dropped \\", size=12.5, color=RGBColor(0xEC,0xEA,0xE4), font=MONO, after=5)
    para(c, "  --tools salesforce gainsight slack", size=12.5, color=RGBColor(0xEC,0xEA,0xE4), font=MONO, after=0)
    say(s, 6.95, 4.85, 5.7, 1.8, "Say it out loud",
        "Three agents run here. The complexity did not disappear, it moved into the platform "
        "where it is written once and tested once.")
    footer(s, page())

    # 12 · Findings --------------------------------------------------------
    s = blank(prs)
    header(s, "11 · Findings from building it", "Three things worth the room's attention")
    items = [
        ("Hallucination is checkable by code",
         "The model cites evidence as metric/value pairs, so every claim is verified against the "
         "data we fetched, deterministically. No second model grading the first. Grounding is 100% "
         "because ungrounded output never survives."),
        ("A single-run eval measures luck",
         "On a deliberately ambiguous case the model gave a different confident answer on a repeat "
         "run. The eval now samples each case several times and gates on consistency, not one lucky run."),
        ("The QA agent has no LLM, deliberately",
         "Owner present, reviews on time, DLQ clean, gate green: each has a correct answer. A model "
         "in the safety path turns a reliable result probabilistic and bills for it."),
    ]
    top = 2.1
    for t, body in items:
        rect(s, 0.7, top, 12, 1.45, PANEL, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, 0.7, top, 0.06, 1.45, ACCENT)
        tf = box(s, 1.0, top + 0.16, 11.4, 1.2, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t, size=16.5, color=INK, bold=True, after=5)
        para(tf, body, size=12.5, color=MUTED, after=0)
        top += 1.62
    footer(s, page())

    # 13 · Failure modes ---------------------------------------------------
    s = blank(prs)
    header(s, "12 · When things go wrong", "The rule: route to a human, never stop silently")
    table(s, 0.7, 2.15, 12, None, [4.4, 7.6],
          ["Failure", "What happens"],
          [
            [("Model invents evidence", {"bold": True}),
             ("Analysis discarded, deterministic answer used, run flagged.", {"color": MUTED})],
            [("Model or vendor down", {"bold": True}),
             ("Fallback chain, then rules. Marked degraded. The human is still alerted.", {"color": MUTED})],
            [("Prediction not trusted", {"bold": True}),
             ("CRM writes held, a person asked. Golden Record marks it awaiting_approval.", {"color": MUTED})],
            [("Flapping trigger loops", {"bold": True}),
             ("Per-account hourly cap on model calls; run continues deterministically.", {"color": MUTED})],
            [("Daily budget runs out", {"bold": True}),
             ("Spending stops at 90%; alerts are never dropped to save money.", {"color": MUTED})],
            [("Two writers, one account", {"bold": True}),
             ("Optimistic concurrency: the stale write is rejected and retried, not lost.", {"color": MUTED})],
          ], size=13)
    footer(s, page())

    # 14 · Honestly --------------------------------------------------------
    s = blank(prs)
    header(s, "13 · Honestly", "What is not real")
    tf = box(s, 0.7, 2.2, 12, 3)
    bullets(tf, [
        ("Vendor clients are mocks ", "over a fixture dataset. They honour the real interface, "
         "retry semantics and idempotency; only the model call crosses the network."),
        ("The BigQuery adapter's SQL is asserted in tests ", "but has never run against a live "
         "dataset: no GCP credentials for this exercise."),
        ("Pseudonymising before the model call is minimisation, ", "not GDPR anonymisation."),
    ], size=17, after=16)
    say(s, 0.7, 5.5, 12, 1.0, "Everything else executes",
        "91 tests, the golden eval gate, and the platform's own self-audit, all green in CI on "
        "three Python versions.")
    footer(s, page())

    # 15 · Close (dark) ----------------------------------------------------
    s = blank(prs, INK)
    rect(s, 0.9, 2.4, 1.05, 0.05, ACCENT)
    tf = box(s, 0.85, 2.6, 11.6, 3)
    para(tf, "Built to be operated, not demoed", size=34, color=WHITE, bold=True, after=16)
    para(tf, "91 tests, an eval gate and a self-audit, all green in CI. What I'd build next: the "
             "approval loop closing back to calibration, async ingestion, and a real BigQuery run.",
         size=18, color=RGBColor(0xC9,0xC7,0xC0), after=0)
    tf2 = box(s, 0.85, 5.7, 11.6, 1)
    runs(tf2, [("runner.py", {"size": 14, "color": RGBColor(0xE0,0x7A,0x5F), "font": MONO}),
               ("  ·  cli.py eval  ·  cli.py audit   —   every claim here is reproducible from the repo.",
                {"size": 14, "color": RGBColor(0x98,0x95,0x8C)})], after=0)

    out = ROOT / "Renewal-Risk-Platform.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    slides = len(Presentation(str(path)).slides._sldIdLst)
    print(f"Wrote {path.name} ({path.stat().st_size // 1024} KB, {slides} slides)")
